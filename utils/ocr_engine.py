# ocr_engine.py
from pathlib import Path
from typing import Optional
import io

# PDF
import fitz  # PyMuPDF

# OCR
import pytesseract
from PIL import Image

# Office formats
import docx  # python-docx
from pptx import Presentation
import openpyxl

# Optional: tune Tesseract binary path if needed (comment / configure in production)
# pytesseract.pytesseract.tesseract_cmd = r"/usr/bin/tesseract"

def extract_text_from_file(path: str, filename: Optional[str] = None) -> str:
    """
    Unified entrypoint the rest of the app can call.

    Args:
      - path: local filesystem path to the file (tempfile created by upload)
      - filename: original uploaded filename (preferred for extension detection)
    Returns:
      - Extracted text (empty string on failure)
    """
    try:
        # Prefer filename extension if provided (keeps original ext when temp file lost ext)
        src = filename if filename else path
        ext = Path(src).suffix.lower().lstrip(".")
        if ext == "pdf":
            return extract_from_pdf(path)
        elif ext in ("jpg", "jpeg", "png", "bmp", "tiff"):
            return extract_from_image(path)
        elif ext in ("docx", "doc"):
            return extract_from_docx(path)
        elif ext == "txt":
            return extract_from_txt(path)
        elif ext in ("xlsx", "xls"):
            return extract_from_xlsx(path)
        elif ext == "pptx":
            return extract_from_pptx(path)
        else:
            return f"[Unsupported file type: .{ext}]"
    except Exception as e:
        return f"[Error reading {filename or path}: {str(e)}]"


# ------------------------- helpers ------------------------- #
def extract_from_pdf(path: str) -> str:
    """
    Extract text from PDF using PyMuPDF (fitz). If a page has no text,
    perform OCR on that page and append OCR result.
    """
    all_text = []
    try:
        doc = fitz.open(path)
    except Exception as e:
        return f"[PDF open error: {e}]"

    for page_index, page in enumerate(doc):
        try:
            # Try structured text extraction first
            page_text = page.get_text("text") or ""
            if page_text.strip():
                all_text.append(page_text.strip())
                continue

            # Fallback OCR: render page to image (higher resolution for better OCR)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x scale
            img_bytes = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_bytes))

            # Preprocess: convert to grayscale (helps OCR often)
            gray = image.convert("L")

            # Optional further binarization could be applied here (adaptive thresholding)
            ocr_text = pytesseract.image_to_string(gray, lang="eng")  # default 'eng' - adjust if using multiple langs
            if ocr_text.strip():
                all_text.append(f"[Page {page_index+1} OCR]\n{ocr_text.strip()}")
        except Exception as e:
            all_text.append(f"[Error on page {page_index+1}: {e}]")

    doc.close()
    return "\n\n".join(all_text).strip()


def extract_from_image(path: str) -> str:
    """Run OCR on an image file"""
    try:
        image = Image.open(path)
        gray = image.convert("L")
        text = pytesseract.image_to_string(gray, lang="eng")
        return text.strip()
    except Exception as e:
        return f"[Image OCR error: {e}]"


def extract_from_docx(path: str) -> str:
    """Extract text from .docx using python-docx"""
    try:
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        return "\n".join(paragraphs).strip()
    except Exception as e:
        return f"[DOCX read error: {e}]"


def extract_from_txt(path: str) -> str:
    """Read plain text file"""
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()
    except Exception as e:
        return f"[TXT read error: {e}]"


def extract_from_xlsx(path: str) -> str:
    """Extract text from Excel using openpyxl (data_only=True)"""
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        out_lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for row in ws.iter_rows(values_only=True):
                # join non-empty cells in row
                cells = [str(c) for c in row if c is not None and str(c).strip() != ""]
                if cells:
                    out_lines.append(" ".join(cells))
        return "\n".join(out_lines).strip()
    except Exception as e:
        return f"[XLSX read error: {e}]"


def extract_from_pptx(path: str) -> str:
    """Extract visible text from PowerPoint slides"""
    try:
        prs = Presentation(path)
        out_lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                # TextFrame shapes (title, body, etc)
                if hasattr(shape, "text") and shape.text:
                    out_lines.append(shape.text.strip())
        return "\n".join(out_lines).strip()
    except Exception as e:
        return f"[PPTX read error: {e}]"
